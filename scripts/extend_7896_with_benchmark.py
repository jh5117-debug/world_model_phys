from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path("/home/hj/Multi-View-Physically-Consistent-World-Model/Physical_Consistency")
SRC_PPT = ROOT / "7896.pptx"
DST_PPT = ROOT / "7896_benchmark_v2.pptx"


def rgb(hex_code: str) -> RGBColor:
    hex_code = hex_code.replace("#", "")
    return RGBColor.from_string(hex_code)


TITLE_COLOR = rgb("#1F2937")
SUBTITLE_COLOR = rgb("#6B7280")
BOX_FILL = rgb("#FBFBFD")
BOX_LINE = rgb("#D9DEE7")
PILL_PURPLE = rgb("#EADCFD")
PILL_BLUE = rgb("#D9EEFF")
PILL_YELLOW = rgb("#FFE9C8")
TEXT_COLOR = rgb("#374151")


def set_text_style(run, *, size: int, bold: bool = False, color: RGBColor = TEXT_COLOR) -> None:
    font = run.font
    font.name = "Microsoft YaHei"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color


def add_title_block(slide, title: str, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_text_style(run, size=28, bold=True, color=TITLE_COLOR)

    subtitle_box = slide.shapes.add_textbox(Inches(0.82), Inches(0.58), Inches(11.0), Inches(0.25))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    set_text_style(run, size=12, color=SUBTITLE_COLOR)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0.86),
        Inches(13.33),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BOX_LINE
    line.line.color.rgb = BOX_LINE


def add_panel(slide, *, x: float, y: float, w: float, h: float, pill_text: str, pill_fill: RGBColor, body_lines: list[str]) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = BOX_FILL
    panel.line.color.rgb = BOX_LINE
    panel.line.width = Pt(1.1)

    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x + 0.18),
        Inches(y + 0.18),
        Inches(2.0),
        Inches(0.34),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = pill_fill
    pill.line.color.rgb = pill_fill

    pill_tf = pill.text_frame
    pill_tf.clear()
    p = pill_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = pill_text
    set_text_style(run, size=11, bold=True, color=TITLE_COLOR)

    body = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.65), Inches(w - 0.5), Inches(h - 0.8))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for idx, line in enumerate(body_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        set_text_style(run, size=14)


def add_footer_note(slide, note: str) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.7),
        Inches(6.85),
        Inches(12.0),
        Inches(0.3),
    )
    box.fill.background()
    box.line.color.rgb = BOX_LINE
    box.line.width = Pt(1.0)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = note
    set_text_style(run, size=11, color=SUBTITLE_COLOR)


def add_workflow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "Workflow Overview", "从 Stage 1 continuation 到统一推理与统一评测")

    add_panel(
        slide,
        x=0.7,
        y=1.25,
        w=6.0,
        h=3.6,
        pill_text="训练优化流程",
        pill_fill=PILL_PURPLE,
        body_lines=[
            "1. 输入 CSGO clip、文本提示和对应条件信息，得到训练视频样本。",
            "2. Teacher 视频编码器提取时空 token 表征 y，提供物理相关关系目标。",
            "3. Student 视频扩散骨干提取中间特征 h，并通过投影层完成维度对齐。",
            "4. 同时优化扩散重建目标 Ldiff 与关系蒸馏目标 LTRD。",
            "5. 分别完成 low 分支与 high 分支优化，最终组成 dual 模型用于正式推理。",
        ],
    )

    add_panel(
        slide,
        x=7.0,
        y=1.25,
        w=5.6,
        h=2.0,
        pill_text="三组对比",
        pill_fill=PILL_BLUE,
        body_lines=[
            "A. LingBot base dual",
            "B. LingBot Stage 1 dual",
            "C. LingBot Stage 1 + VideoREPA-style TRD dual",
        ],
    )

    add_panel(
        slide,
        x=7.0,
        y=3.45,
        w=5.6,
        h=1.95,
        pill_text="统一推理设置",
        pill_fill=PILL_YELLOW,
        body_lines=[
            "同一 benchmark split、同一 prompt/首帧/控制信号、同一采样步数、同一 seed。",
            "这样最终差异可以尽量归因到模型本身，而不是推理条件差异。",
        ],
    )

    add_footer_note(
        slide,
        "一句话：训练阶段用 TRD 注入时空关系约束，推理阶段统一设定，评测阶段统一协议，形成闭环。",
    )


def add_benchmark_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "VideoPhy-2 AutoEval Benchmark", "当前无独立 test split，因此使用 processed_csgo_v3/val 作为 internal benchmark")

    add_panel(
        slide,
        x=0.7,
        y=1.2,
        w=6.0,
        h=2.35,
        pill_text="Benchmark 输入",
        pill_fill=PILL_BLUE,
        body_lines=[
            "评测 split：/home/nvme02/lingbot-world/datasets/processed_csgo_v3/val/",
            "比较对象：LingBot(base dual)、LingBot(Stage 1 dual)、LingBot(Stage 1 + VideoREPA-style TRD dual)。",
            "对 val 中每个 clip，在相同条件下生成视频并送入 AutoEval。",
        ],
    )

    add_panel(
        slide,
        x=6.95,
        y=1.2,
        w=5.7,
        h=2.35,
        pill_text="AutoEval 任务",
        pill_fill=PILL_PURPLE,
        body_lines=[
            "SA：Semantic Adherence，输入 videopath + caption，评估视频是否符合文本语义。",
            "PC：Physical Commonsense，输入 videopath，评估视频是否符合物理常识。",
            "每个样本输出 1 到 5 分。",
        ],
    )

    add_panel(
        slide,
        x=0.7,
        y=3.85,
        w=11.95,
        h=2.25,
        pill_text="汇总指标",
        pill_fill=PILL_YELLOW,
        body_lines=[
            "SA_mean = (1/N) Σ_i s_i^SA",
            "PC_mean = (1/N) Σ_i s_i^PC",
            "Joint = (1/N) Σ_i 1[(s_i^SA ≥ 4) ∧ (s_i^PC ≥ 4)]",
            "最终报告时同时给出 SA、PC 和 Joint，并明确说明这是当前的 internal benchmark。",
        ],
    )

    add_footer_note(
        slide,
        "一句话：先统一生成，再统一做 VideoPhy-2 AutoEval，最后用 SA / PC / Joint 比较三组模型的物理一致性表现。",
    )


def main() -> None:
    prs = Presentation(str(SRC_PPT))
    add_workflow_slide(prs)
    add_benchmark_slide(prs)
    prs.save(str(DST_PPT))
    print(DST_PPT)


if __name__ == "__main__":
    main()
