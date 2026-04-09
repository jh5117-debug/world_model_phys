#!/usr/bin/env python3
"""Generate a PowerPoint deck from the pre-PPT markdown document."""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

COLOR_BG_DARK = RGBColor(14, 23, 38)
COLOR_BG_LIGHT = RGBColor(248, 250, 252)
COLOR_TEXT = RGBColor(15, 23, 42)
COLOR_MUTED = RGBColor(71, 85, 105)
COLOR_ACCENT = RGBColor(249, 115, 22)
COLOR_ACCENT_LIGHT = RGBColor(255, 247, 237)
COLOR_BORDER = RGBColor(226, 232, 240)

FONT_FAMILY = "Microsoft YaHei"

LABELS = [
    "页面标题",
    "页面内容",
    "关键说明",
    "当前目录",
    "关键文件",
    "关键代码",
    "建议图示",
    "扩展讲稿",
]


@dataclass
class BulletLine:
    level: int
    text: str
    kind: str = "bullet"


@dataclass
class SlideSpec:
    heading: str
    title: str
    body: list[BulletLine] = field(default_factory=list)
    extra: list[tuple[str, list[BulletLine]]] = field(default_factory=list)
    notes: str = ""
    is_title: bool = False
    is_appendix: bool = False


def split_sections(text: str) -> list[tuple[str, str]]:
    pattern = re.compile(r"^##\s+(.+)$", re.MULTILINE)
    matches = list(pattern.finditer(text))
    sections: list[tuple[str, str]] = []
    for idx, match in enumerate(matches):
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        heading = match.group(1).strip()
        body = text[start:end].strip()
        sections.append((heading, body))
    return sections


def extract_label_blocks(body: str) -> dict[str, str]:
    blocks: dict[str, list[str]] = {}
    current: str | None = None
    label_re = re.compile(rf"^({'|'.join(map(re.escape, LABELS))})：\s*$")
    for raw_line in body.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        match = label_re.match(stripped)
        if match:
            current = match.group(1)
            blocks.setdefault(current, [])
            continue
        if current is not None:
            blocks[current].append(line)
    return {key: "\n".join(value).strip() for key, value in blocks.items()}


def parse_lines(block: str) -> list[BulletLine]:
    if not block:
        return []

    lines: list[BulletLine] = []
    for raw in block.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue

        bullet_match = re.match(r"^(\s*)-\s+(.*)$", line)
        number_match = re.match(r"^(\s*)(\d+)\.\s+(.*)$", line)
        if bullet_match:
            indent = len(bullet_match.group(1).replace("\t", "    ")) // 2
            lines.append(BulletLine(level=indent, text=clean_inline(bullet_match.group(2)), kind="bullet"))
            continue
        if number_match:
            indent = len(number_match.group(1).replace("\t", "    ")) // 2
            number = number_match.group(2)
            text = clean_inline(number_match.group(3))
            lines.append(BulletLine(level=indent, text=f"{number}. {text}", kind="number"))
            continue

        text = clean_inline(line.strip())
        if lines:
            lines[-1].text = f"{lines[-1].text} {text}"
        else:
            lines.append(BulletLine(level=0, text=text, kind="bullet"))
    return lines


def clean_inline(text: str) -> str:
    text = text.strip()
    if text.startswith("`") and text.endswith("`") and len(text) >= 2:
        text = text[1:-1]
    return text


def build_slide_specs(md_path: Path) -> list[SlideSpec]:
    text = md_path.read_text(encoding="utf-8")
    sections = split_sections(text)
    slides: list[SlideSpec] = []

    for heading, body in sections:
        if not (heading.startswith("第 ") or heading.startswith("附录 ")):
            continue

        if heading.startswith("附录 "):
            appendix_lines = parse_lines(body)
            slides.append(
                SlideSpec(
                    heading=heading,
                    title=heading,
                    body=appendix_lines,
                    is_appendix=True,
                )
            )
            continue

        blocks = extract_label_blocks(body)
        title = clean_inline(blocks.get("页面标题", heading))
        body_lines = parse_lines(blocks.get("页面内容", ""))
        extra_sections: list[tuple[str, list[BulletLine]]] = []
        for key in ["关键说明", "当前目录", "关键文件", "关键代码"]:
            if key in blocks and blocks[key].strip():
                extra_sections.append((key, parse_lines(blocks[key])))

        notes = blocks.get("扩展讲稿", "").strip()
        slides.append(
            SlideSpec(
                heading=heading,
                title=title,
                body=body_lines,
                extra=extra_sections,
                notes=notes,
                is_title=heading.startswith("第 1 页"),
            )
        )
    return slides


def add_full_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_textbox(slide, left, top, width, height, text="", *, font_size=18, bold=False, color=COLOR_TEXT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = FONT_FAMILY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return textbox


def style_paragraph(paragraph, *, font_size=18, color=COLOR_TEXT, bold=False, level=0, align=PP_ALIGN.LEFT):
    paragraph.level = max(level, 0)
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def estimate_body_font(lines: list[BulletLine]) -> int:
    count = len(lines)
    total_chars = sum(len(line.text) for line in lines)
    if count >= 9 or total_chars > 520:
        return 15
    if count >= 7 or total_chars > 360:
        return 16
    return 18


def estimate_extra_font(lines: list[tuple[str, list[BulletLine]]]) -> int:
    total_items = sum(len(items) + 1 for _, items in lines)
    total_chars = sum(len(item.text) for _, items in lines for item in items)
    if total_items > 12 or total_chars > 500:
        return 9
    if total_items > 8 or total_chars > 320:
        return 10
    return 11


def add_footer(slide, slide_no: int):
    add_textbox(
        slide,
        Inches(12.2),
        Inches(7.0),
        Inches(0.8),
        Inches(0.25),
        str(slide_no),
        font_size=10,
        bold=False,
        color=COLOR_MUTED,
    )


def build_title_slide(slide, spec: SlideSpec, slide_no: int):
    add_full_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_BG_DARK)
    add_full_rect(slide, Inches(0.6), Inches(0.72), Inches(0.14), Inches(5.2), COLOR_ACCENT)
    add_textbox(
        slide,
        Inches(1.05),
        Inches(0.75),
        Inches(10.6),
        Inches(1.6),
        spec.title,
        font_size=28,
        bold=True,
        color=RGBColor(255, 255, 255),
    )

    subtitle_lines = [line.text for line in spec.body]
    subtitle = "\n".join(f"• {line}" for line in subtitle_lines[:4])
    sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.45), Inches(6.4), Inches(2.2))
    frame = sub_box.text_frame
    frame.word_wrap = True
    for idx, item in enumerate(subtitle.splitlines()):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=18, color=RGBColor(241, 245, 249), level=0)

    callout = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.35),
        Inches(2.4),
        Inches(4.1),
        Inches(2.5),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(30, 41, 59)
    callout.line.color.rgb = COLOR_ACCENT
    tf = callout.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(["Physics", "Evaluation", "Isolated Engineering"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=20, color=RGBColor(255, 247, 237), bold=True)

    add_textbox(
        slide,
        Inches(1.1),
        Inches(6.35),
        Inches(5.2),
        Inches(0.45),
        "Physical_Consistency / Stage 1 continuation / CSGO",
        font_size=12,
        color=RGBColor(203, 213, 225),
    )
    add_footer(slide, slide_no)


def build_content_slide(slide, spec: SlideSpec, slide_no: int):
    add_full_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, RGBColor(255, 255, 255))
    add_full_rect(slide, Inches(0.0), Inches(0.0), SLIDE_WIDTH, Inches(0.18), COLOR_ACCENT)

    add_textbox(slide, Inches(0.6), Inches(0.38), Inches(10.8), Inches(0.5), spec.title, font_size=24, bold=True)
    add_textbox(
        slide,
        Inches(11.65),
        Inches(0.42),
        Inches(1.0),
        Inches(0.35),
        spec.heading.replace("第 ", "P").replace(" 页：", " "),
        font_size=10,
        color=COLOR_MUTED,
    )

    has_extra = bool(spec.extra)
    body_left = Inches(0.65)
    body_top = Inches(1.05)
    body_width = Inches(7.7 if has_extra else 11.9)
    body_height = Inches(5.75)

    body_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        body_left,
        body_top,
        body_width,
        body_height,
    )
    body_box.fill.solid()
    body_box.fill.fore_color.rgb = COLOR_BG_LIGHT
    body_box.line.color.rgb = COLOR_BORDER
    body_box.line.width = Pt(1)

    frame = body_box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.15)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)

    font_size = estimate_body_font(spec.body)
    if spec.body:
        for idx, item in enumerate(spec.body):
            p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            prefix = "• " if not item.text[:2].isdigit() else ""
            p.text = f"{prefix}{item.text}"
            style_paragraph(p, font_size=font_size, level=item.level)
    else:
        p = frame.paragraphs[0]
        p.text = " "

    if has_extra:
        extra_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.7),
            Inches(1.05),
            Inches(3.95),
            Inches(5.75),
        )
        extra_box.fill.solid()
        extra_box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
        extra_box.line.color.rgb = RGBColor(253, 186, 116)
        extra_box.line.width = Pt(1)
        ef = extra_box.text_frame
        ef.word_wrap = True
        ef.margin_left = Inches(0.16)
        ef.margin_right = Inches(0.12)
        ef.margin_top = Inches(0.12)
        ef.margin_bottom = Inches(0.08)
        extra_font = estimate_extra_font(spec.extra)
        first = True
        for label, items in spec.extra:
            hp = ef.paragraphs[0] if first else ef.add_paragraph()
            hp.text = label
            style_paragraph(hp, font_size=12, color=COLOR_ACCENT, bold=True)
            first = False
            for item in items:
                p = ef.add_paragraph()
                prefix = "• " if not item.text[:2].isdigit() else ""
                p.text = f"{prefix}{item.text}"
                style_paragraph(p, font_size=extra_font, level=min(item.level, 1))

    if spec.notes:
        note_text = spec.notes.strip().replace("\n", " ")
        if len(note_text) > 170:
            note_text = note_text[:167] + "..."
        note_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.7),
            Inches(6.95),
            Inches(10.8),
            Inches(0.3),
        )
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
        note_box.line.color.rgb = COLOR_BORDER
        ntf = note_box.text_frame
        ntf.word_wrap = True
        p = ntf.paragraphs[0]
        p.text = f"讲稿提示：{note_text}"
        style_paragraph(p, font_size=9, color=COLOR_MUTED)

    add_footer(slide, slide_no)


def build_appendix_slide(slide, spec: SlideSpec, slide_no: int):
    add_full_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, RGBColor(255, 255, 255))
    add_full_rect(slide, Inches(0.0), Inches(0.0), SLIDE_WIDTH, Inches(0.18), COLOR_ACCENT)
    add_textbox(slide, Inches(0.6), Inches(0.38), Inches(11.6), Inches(0.5), spec.title, font_size=22, bold=True)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.0),
        Inches(12.0),
        Inches(5.9),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_BG_LIGHT
    box.line.color.rgb = COLOR_BORDER
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)

    font_size = 13 if len(spec.body) > 12 else 15
    for idx, item in enumerate(spec.body):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        prefix = "• " if not item.text[:2].isdigit() else ""
        p.text = f"{prefix}{item.text}"
        style_paragraph(p, font_size=font_size, level=min(item.level, 2))

    add_footer(slide, slide_no)


def generate_ppt(md_path: Path, out_path: Path) -> Path:
    slides = build_slide_specs(md_path)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        if spec.is_title:
            build_title_slide(slide, spec, idx)
        elif spec.is_appendix:
            build_appendix_slide(slide, spec, idx)
        else:
            build_content_slide(slide, spec, idx)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate PowerPoint from markdown pre file.")
    parser.add_argument(
        "--input",
        type=str,
        default="/home/hj/Multi-View-Physically-Consistent-World-Model/Physical_Consistency/PPT_physical_consistency_pre_v1.md",
    )
    parser.add_argument(
        "--output",
        type=str,
        default="/home/hj/Multi-View-Physically-Consistent-World-Model/Physical_Consistency/PPT_physical_consistency_pre_v1.pptx",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = generate_ppt(Path(args.input), Path(args.output))
    print(output)


if __name__ == "__main__":
    main()
