#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_PATH = ROOT / "WEEKLY_REPORT_physical_consistency_2026-04-14.pptx"
SCRIPT_PATH = ROOT / "WEEKLY_REPORT_physical_consistency_2026-04-14_汇报稿.md"


SLIDES = [
    {
        "title": "本周工作目标与总体进展",
        "subtitle": "先把物理一致性训练与测试的基础系统搭起来",
        "bullets": [
            "本周重点完成了测试流程、benchmark 流程、训练框架、验证框架和工程 debug。",
            "当前已经形成一套可训练、可验证、可持续迭代的基础链路。",
            "这周的核心目标不是追最终指标，而是先把系统真正跑通。",
        ],
        "takeaway": "阶段结论：系统已经从零散代码推进到可用框架。",
        "script": [
            "这周我的重点，不是直接追最终效果，而是先把物理一致性这一块的基础系统搭起来。",
            "我主要完成了测试流程、benchmark 流程、训练框架、验证框架，以及一系列关键的工程 debug。",
            "到目前为止，这套系统已经从零散代码，变成了一个可以稳定训练、可以接验证、也可以持续优化的基础链路。",
        ],
    },
    {
        "title": "测试与 Benchmark 搭建",
        "subtitle": "把训练、生成和评测之间的关系先理顺",
        "bullets": [
            "我梳理了训练集、验证集和评测入口，明确了每一步数据流向。",
            "我把生成结果到评测结果的基本流程打通，方便后续统一比较实验。",
            "现在支持按 epoch 自动触发验证，训练和评测形成了基本闭环。",
        ],
        "takeaway": "阶段结论：后续改模型时，不再只看 loss，而是可以稳定接到生成和评测。",
        "script": [
            "这周我先把测试和 benchmark 的基础流程理清了。",
            "我明确了训练数据、验证数据、生成结果和评测结果之间的对应关系，也把验证流程接到了训练过程中。",
            "现在每个 epoch 结束后，系统可以自动进入验证流程，这样后面我们做模型改动时，就不只是看 loss，而是能稳定地接到生成和评测。",
        ],
    },
    {
        "title": "物理一致性训练框架搭建",
        "subtitle": "保留稳定主框架，只替换必要模块",
        "bullets": [
            "我把 Stage1 + TRD dual student 的训练主流程真正串了起来。",
            "我尽量保留原来已经稳定的框架，只在必要位置做替换，降低系统风险。",
            "当前 teacher 已经从 VideoMAEv2 切换到官方 V-JEPA 2.1，并能接入训练主循环。",
        ],
        "takeaway": "阶段结论：训练主链路已经稳定可跑，teacher 替换也已落地。",
        "script": [
            "在训练框架上，我这周重点是把物理一致性的训练主流程真正串起来。",
            "我采用的是保守策略，就是尽量保留原来已经能工作的稳定框架，只在必要位置做替换，这样可以减少系统性风险。",
            "在这个基础上，我已经把 teacher 从原来的 VideoMAEv2 切换到了官方 V-JEPA 2.1，并且训练主循环已经能够正常运行。",
        ],
    },
    {
        "title": "关键 Debug 工作",
        "subtitle": "这周投入最多时间的部分是把关键工程问题逐个定位下来",
        "bullets": [
            "我解决了 OOM、日志混乱、W&B 图表异常、后台训练易中断等问题。",
            "我解决了 nohup 和 SSH 断开导致训练被挂掉的问题，增强了后台启动稳定性。",
            "我定位并修复了 epoch-end validation 挂起问题，根因是验证子进程错误继承了分布式环境。",
        ],
        "takeaway": "阶段结论：最影响实验连续性的工程问题已经基本清掉。",
        "script": [
            "这周花时间最多的部分，其实是 debug。",
            "前面遇到的问题比较多，包括显存 OOM、日志输出过于混乱、W&B 图表异常、后台训练容易被终端状态影响，以及每个 epoch 结束后的验证流程会挂住。",
            "我最后把这些问题逐步定位下来，尤其是 validation 挂起这个问题，根因不是模型本身，而是验证子进程错误继承了训练时的分布式环境，导致它卡在 barrier。我已经把这个问题修复了。",
        ],
    },
    {
        "title": "当前阶段结果",
        "subtitle": "先看链路是否闭环，再看后续指标是否可持续提升",
        "bullets": [
            "训练已经可以稳定启动并持续运行，teacher 新旧切换也已经跑通。",
            "每个 epoch 结束后都可以进入验证流程，训练、生成和评测已经基本接上。",
            "当前最重要的结果，是我们已经有了一个能持续迭代的实验平台。",
        ],
        "takeaway": "阶段结论：现在已经具备“训练 - 生成 - 评测”的基础闭环。",
        "script": [
            "从当前阶段的结果来看，我觉得最重要的不是某一个具体数值，而是整条链路已经基本闭环了。",
            "现在训练已经可以稳定跑起来，teacher 也已经切换到新的官方模型，同时每个 epoch 结束后都可以接上验证流程。",
            "这说明我们现在已经有了一个可以持续做实验的平台，后面不管是改 teacher、改 loss，还是加入新的物理约束，都有了比较稳定的基础。",
        ],
    },
    {
        "title": "下周计划与核心创新点",
        "subtitle": "在稳定框架上，开始更系统地做实验与总结",
        "bullets": [
            "我会基于当前稳定框架，继续做更系统的对比实验。",
            "重点分析 teacher 变化、物理一致性 loss 和不同设计选择分别带来的贡献。",
            "目标是把我们的核心创新点提炼得更清楚，不只是系统能跑，而是方法真正新在哪里、强在哪里。",
        ],
        "takeaway": "阶段结论：下周重点从“修系统”转向“提炼创新点”。",
        "script": [
            "下周我的重点就不再只是修工程问题，而是基于这套已经稳定的框架，开始更系统地做实验。",
            "我会重点分析 teacher 变化、物理一致性 loss，以及不同设计选择分别带来了什么贡献。",
            "我希望下周能够更清楚地提炼出我们的方法到底新在哪里、强在哪里，也就是把我们的核心创新点讲得更明确。",
        ],
    },
]


def set_run_style(run, *, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    font = run.font
    font.name = "Microsoft YaHei"
    font.size = Pt(size)
    font.bold = bold
    if color is not None:
        font.color.rgb = color


def add_text_frame(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    paragraphs: list[str],
    font_size: int,
    color: RGBColor,
    bullet: bool = False,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for idx, text in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(8)
        if bullet:
            p.bullet = True
        for run in p.runs:
            set_run_style(run, size=font_size, color=color)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(248, 250, 252)
    title_color = RGBColor(17, 24, 39)
    body_color = RGBColor(55, 65, 81)
    accent_color = RGBColor(14, 116, 144)
    highlight_fill = RGBColor(224, 242, 254)
    highlight_text = RGBColor(12, 74, 110)

    for index, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = bg_color

        top_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = accent_color
        top_bar.line.fill.background()

        slide_num = slide.shapes.add_textbox(Inches(12.2), Inches(0.28), Inches(0.8), Inches(0.3))
        tf = slide_num.text_frame
        p = tf.paragraphs[0]
        p.text = f"{index:02d}"
        p.alignment = PP_ALIGN.RIGHT
        for run in p.runs:
            set_run_style(run, size=14, bold=True, color=accent_color)

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(10.8), Inches(0.7))
        p = title_box.text_frame.paragraphs[0]
        p.text = spec["title"]
        for run in p.runs:
            set_run_style(run, size=24, bold=True, color=title_color)

        subtitle_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.18), Inches(11.2), Inches(0.5))
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = spec["subtitle"]
        for run in p.runs:
            set_run_style(run, size=12, color=accent_color)

        add_text_frame(
            slide,
            left=0.8,
            top=1.8,
            width=11.6,
            height=3.4,
            paragraphs=spec["bullets"],
            font_size=20,
            color=body_color,
            bullet=True,
        )

        takeaway = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.7), Inches(11.85), Inches(0.9)
        )
        takeaway.fill.solid()
        takeaway.fill.fore_color.rgb = highlight_fill
        takeaway.line.color.rgb = RGBColor(186, 230, 253)

        tf = takeaway.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = spec["takeaway"]
        for run in p.runs:
            set_run_style(run, size=16, bold=True, color=highlight_text)

        footer = slide.shapes.add_textbox(Inches(0.75), Inches(6.9), Inches(5.0), Inches(0.25))
        p = footer.text_frame.paragraphs[0]
        p.text = "Physical Consistency Weekly Update"
        for run in p.runs:
            set_run_style(run, size=10, color=RGBColor(107, 114, 128))

    return prs


def build_script_markdown() -> str:
    lines = [
        "# Physical Consistency 周汇报讲稿",
        "",
        "下面按 PPT 页面顺序给出逐页汇报稿。",
        "",
    ]
    for idx, slide in enumerate(SLIDES, start=1):
        lines.append(f"## 第{idx}页：{slide['title']}")
        lines.append("")
        for paragraph in slide["script"]:
            lines.append(paragraph)
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def main() -> None:
    prs = build_presentation()
    prs.save(PPT_PATH)
    SCRIPT_PATH.write_text(build_script_markdown(), encoding="utf-8")
    print(PPT_PATH)
    print(SCRIPT_PATH)


if __name__ == "__main__":
    main()
